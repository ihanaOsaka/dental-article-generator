#!/usr/bin/env python3
"""Generate PowerPoint (.pptx) slide decks from video scenario markdown files.

Output .pptx files can also be imported directly into Google Slides.
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SCENARIOS_DIR = Path("output/video_scenarios")
PPTX_DIR = Path("output/slides_pptx")

# Color scheme
COLOR_PRIMARY = RGBColor(0x1B, 0x5E, 0x7B)    # Deep teal
COLOR_ACCENT = RGBColor(0x26, 0xA6, 0x9A)     # Teal accent
COLOR_BG_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)   # Light teal bg
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)        # Dark gray text
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SUBTITLE = RGBColor(0x55, 0x55, 0x55)
COLOR_QUESTION_BG = RGBColor(0xF5, 0xF5, 0xF5)


def set_slide_bg(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_fill(slide, left, top, width, height, color):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=COLOR_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Meiryo"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=COLOR_TEXT, font_name="Meiryo"):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
        # Add bullet character
        p.text = f"• {item}"

    return txBox


def parse_scenario(filepath: Path) -> dict:
    """Parse a video scenario markdown file into structured data."""
    text = filepath.read_text(encoding="utf-8")

    title_match = re.search(r"^# 動画シナリオ #(\d+)：(.+)$", text, re.MULTILINE)
    number = title_match.group(1)
    title = title_match.group(2)

    audience_match = re.search(r"\| 対象視聴者 \| (.+?) \|", text)
    audience = audience_match.group(1) if audience_match else ""

    key_msg_match = re.search(r"\*\*キーメッセージ:\*\* (.+?)$", text, re.MULTILINE)
    key_message = key_msg_match.group(1).strip().strip("「」") if key_msg_match else ""

    duration_match = re.search(r"\| 想定尺 \| (.+?) \|", text)
    duration = duration_match.group(1) if duration_match else ""

    sections = []
    section_pattern = re.compile(
        r"### (.+?)(?:\n\n|\n)"
        r"(.*?)(?=\n### |\n## 制作メモ|\Z)",
        re.DOTALL,
    )
    for match in section_pattern.finditer(text):
        header = match.group(1).strip()
        body = match.group(2).strip()

        telop_match = re.search(r"\*\*テロップ:\*\* (.+?)$", body, re.MULTILINE)
        telop = telop_match.group(1).strip() if telop_match else ""

        narration = ""
        nar_match = re.search(
            r"\*\*ナレーション:\*\*\n(.+?)(?=\n\*\*|$)", body, re.DOTALL
        )
        if nar_match:
            narration = nar_match.group(1).strip()
            narration = re.sub(r"^「|」$", "", narration.strip())

        question = ""
        q_match = re.search(
            r"\*\*質問（画面テキスト\+読み上げ）:\*\*\n(.+?)(?=\n\*\*|$)",
            body,
            re.DOTALL,
        )
        if q_match:
            question = q_match.group(1).strip()
            question = re.sub(r"^「|」$", "", question.strip())

        visual = ""
        vis_match = re.search(
            r"\*\*画面演出案:\*\*\n(.+?)(?=\n---|\Z)", body, re.DOTALL
        )
        if vis_match:
            visual = vis_match.group(1).strip()
            # Clean up markdown list markers
            visual = re.sub(r"^- ", "", visual, flags=re.MULTILINE)

        sections.append({
            "header": header,
            "telop": telop,
            "narration": narration,
            "question": question,
            "visual": visual,
        })

    return {
        "number": number,
        "title": title,
        "audience": audience,
        "duration": duration,
        "key_message": key_message,
        "sections": sections,
    }


def extract_narration_bullets(narration: str) -> list[str]:
    """Extract key points from narration as bullet points."""
    bullets = []
    lines = narration.split("\n")
    current = []
    for line in lines:
        line = line.strip().strip("「」")
        if not line:
            if current:
                bullets.append("".join(current))
                current = []
        else:
            current.append(line)
    if current:
        bullets.append("".join(current))
    return [b for b in bullets if len(b) > 5]


def create_title_slide(prs, scenario):
    """Create the title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, COLOR_PRIMARY)

    # Top bar accent
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(10), Inches(0.08), COLOR_ACCENT)

    # Number badge
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.6),
                 f"#{scenario['number']}", font_size=20, color=COLOR_ACCENT,
                 bold=True)

    # Title
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5),
                 scenario["title"], font_size=32, color=COLOR_WHITE, bold=True)

    # Audience & duration
    info_text = f"対象: {scenario['audience']}"
    if scenario["duration"]:
        info_text += f"  |  {scenario['duration']}"
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.5),
                 info_text, font_size=14, color=RGBColor(0xB0, 0xD0, 0xD8))

    # Bottom accent line
    add_shape_with_fill(slide, Inches(0.8), Inches(4.3), Inches(2), Inches(0.04), COLOR_ACCENT)


def create_opening_slide(prs, section, scenario):
    """Create the opening/question slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)

    # Top bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), COLOR_PRIMARY)

    # Section label
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.4),
                 "オープニング", font_size=12, color=COLOR_ACCENT, bold=True)

    # Telop as heading
    if section["telop"]:
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8.4), Inches(0.8),
                     section["telop"], font_size=28, color=COLOR_PRIMARY, bold=True)

    # Question box
    if section["question"]:
        q_bg = add_shape_with_fill(slide, Inches(0.6), Inches(2.0), Inches(8.8), Inches(2.5),
                                    COLOR_BG_LIGHT)
        # Left accent bar for quote style
        add_shape_with_fill(slide, Inches(0.6), Inches(2.0), Inches(0.08), Inches(2.5),
                            COLOR_ACCENT)

        add_text_box(slide, Inches(1.0), Inches(2.1), Inches(8.2), Inches(0.4),
                     "視聴者からの質問", font_size=12, color=COLOR_ACCENT, bold=True)

        add_text_box(slide, Inches(1.0), Inches(2.5), Inches(8.2), Inches(1.8),
                     section["question"], font_size=18, color=COLOR_TEXT)


def create_content_slide(prs, section, slide_num):
    """Create a content slide for a section."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)

    # Top bar
    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), COLOR_PRIMARY)

    # Clean header - remove time codes
    header = re.sub(r"（[\d:〜]+）.*$", "", section["header"])

    # Section label
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.4),
                 header, font_size=12, color=COLOR_ACCENT, bold=True)

    # Telop as main heading
    if section["telop"]:
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8.4), Inches(0.8),
                     section["telop"], font_size=24, color=COLOR_PRIMARY, bold=True)

    # Narration bullets
    if section["narration"]:
        bullets = extract_narration_bullets(section["narration"])
        if bullets:
            # Limit bullet text length for readability
            trimmed = []
            for b in bullets:
                if len(b) > 80:
                    trimmed.append(b[:80] + "…")
                else:
                    trimmed.append(b)
            add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(8.4), Inches(2.8),
                           trimmed, font_size=14)

    # Visual notes at bottom
    if section["visual"]:
        vis_lines = section["visual"].split("\n")
        vis_text = " / ".join(l.strip() for l in vis_lines if l.strip())
        add_shape_with_fill(slide, Inches(0.6), Inches(4.6), Inches(8.8), Inches(0.7),
                            COLOR_BG_LIGHT)
        add_text_box(slide, Inches(0.8), Inches(4.65), Inches(8.4), Inches(0.6),
                     f"演出: {vis_text}", font_size=11, color=COLOR_SUBTITLE)


def create_ending_slide(prs, section, scenario):
    """Create the ending slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_PRIMARY)

    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(10), Inches(0.08), COLOR_ACCENT)

    # Telop
    if section["telop"]:
        add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.8),
                     section["telop"], font_size=28, color=COLOR_WHITE, bold=True)

    # Narration
    if section["narration"]:
        bullets = extract_narration_bullets(section["narration"])
        if bullets:
            text = "\n".join(bullets)
            add_text_box(slide, Inches(0.8), Inches(2.3), Inches(8.4), Inches(1.5),
                         text, font_size=16, color=RGBColor(0xCC, 0xE8, 0xED))

    # CTA
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.5),
                 "クリニック情報・ご予約はこちら", font_size=14,
                 color=COLOR_ACCENT, bold=True)


def create_summary_slide(prs, scenario):
    """Create the summary/key-message slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_WHITE)

    add_shape_with_fill(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), COLOR_PRIMARY)

    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(8.4), Inches(0.6),
                 "まとめ", font_size=14, color=COLOR_ACCENT, bold=True)

    # Key message box
    if scenario["key_message"]:
        add_shape_with_fill(slide, Inches(0.6), Inches(1.6), Inches(8.8), Inches(2.0),
                            COLOR_BG_LIGHT)
        add_shape_with_fill(slide, Inches(0.6), Inches(1.6), Inches(0.08), Inches(2.0),
                            COLOR_ACCENT)

        add_text_box(slide, Inches(1.0), Inches(1.7), Inches(8.2), Inches(0.4),
                     "キーメッセージ", font_size=12, color=COLOR_ACCENT, bold=True)
        add_text_box(slide, Inches(1.0), Inches(2.1), Inches(8.2), Inches(1.3),
                     scenario["key_message"], font_size=20, color=COLOR_PRIMARY, bold=True)

    # Contact info placeholder
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.5),
                 "お問い合わせ: クリニック情報・ご予約方法", font_size=14,
                 color=COLOR_SUBTITLE)


def generate_pptx(scenario: dict, output_path: Path):
    """Generate a full .pptx presentation from a parsed scenario."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    sections = scenario["sections"]

    # Slide 1: Title
    create_title_slide(prs, scenario)

    for sec in sections:
        header = sec["header"]
        if "オープニング" in header:
            create_opening_slide(prs, sec, scenario)
        elif "エンディング" in header:
            create_ending_slide(prs, sec, scenario)
        else:
            create_content_slide(prs, sec, sec["header"])

    # Summary slide
    create_summary_slide(prs, scenario)

    prs.save(str(output_path))


def main():
    PPTX_DIR.mkdir(parents=True, exist_ok=True)

    scenario_files = sorted(SCENARIOS_DIR.glob("*.md"))
    print(f"Found {len(scenario_files)} scenarios")

    for filepath in scenario_files:
        scenario = parse_scenario(filepath)
        output_path = PPTX_DIR / filepath.name.replace(".md", ".pptx")
        generate_pptx(scenario, output_path)
        print(f"  Generated: {output_path.name}")

    print(f"\nDone! {len(scenario_files)} .pptx files created in {PPTX_DIR}/")


if __name__ == "__main__":
    main()
